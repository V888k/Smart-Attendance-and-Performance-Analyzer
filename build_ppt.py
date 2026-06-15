"""Build the SAPA project presentation from scratch."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- Presentation setup ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------- Palette ----------
BG = RGBColor(0x0E, 0x14, 0x2B)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
PANEL = RGBColor(0x16, 0x21, 0x3E)
TEAL = RGBColor(0x0F, 0x34, 0x60)
CORAL = RGBColor(0xE9, 0x45, 0x60)
MINT = RGBColor(0x00, 0xB8, 0x94)
AMBER = RGBColor(0xFD, 0xCB, 0x6E)
CRIMSON = RGBColor(0xD6, 0x30, 0x31)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0xE6, 0xE6, 0xF0)
MUTED = RGBColor(0x9A, 0x9A, 0xB8)
DIM = RGBColor(0x66, 0x6A, 0x85)
FONT = "Calibri"

# ---------- Primitive helpers ----------
def bg(slide, color=BG):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def band(slide, l, t, w, h, color, alpha_line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background() if not alpha_line else None
    return s

def panel(slide, l, t, w, h, color=PANEL, corner=None):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.color.rgb = RGBColor(0x2A, 0x33, 0x55)
    s.line.width = Pt(0.5)
    return s

def tb(slide, l, t, w, h, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT
        p.alignment = align
    return box

def page_no(slide, n, total):
    tb(slide, Inches(12.4), Inches(7.15), Inches(0.85), Inches(0.3),
       f"{n} / {total}", size=11, color=DIM, align=PP_ALIGN.RIGHT)

def brand(slide, x=Inches(0.5), y=Inches(7.15)):
    tb(slide, x, y, Inches(4), Inches(0.3),
       "SAPA — Smart Attendance & Performance Analyzer",
       size=11, color=DIM)

# ---------- Slide chrome ----------
def chrome(slide, title, kicker=None, accent=CORAL, total=14, n=0):
    bg(slide, BG)
    # top accent
    band(slide, 0, 0, prs.slide_width, Inches(0.08), accent)
    # title bar
    band(slide, 0, Inches(0.08), prs.slide_width, Inches(1.1), NAVY)
    # title text
    tb(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.55),
       title, size=30, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        tb(slide, Inches(0.6), Inches(0.74), Inches(11.5), Inches(0.35),
           kicker, size=14, color=MUTED)
    # accent under title
    band(slide, Inches(0.6), Inches(1.22), Inches(0.6), Inches(0.05), accent)
    brand(slide)
    if n:
        page_no(slide, n, total)

# ---------- Card factory ----------
def card(slide, l, t, w, h, title, lines, accent=CORAL, body=14, head=18):
    panel(slide, l, t, w, h, PANEL, corner=True)
    # accent dot
    band(slide, l + Inches(0.25), t + Inches(0.28), Inches(0.12), Inches(0.12), accent)
    tb(slide, l + Inches(0.5), t + Inches(0.18), w - Inches(0.7), Inches(0.4),
       title, size=head, color=WHITE, bold=True)
    if isinstance(lines, str):
        lines = [lines]
    tb(slide, l + Inches(0.3), t + Inches(0.7), w - Inches(0.6), h - Inches(0.9),
       "\n".join(["•  " + s for s in lines]),
       size=body, color=TEXT)

# ---------- Stat block ----------
def stat(slide, l, t, w, h, value, label, color=CORAL):
    panel(slide, l, t, w, h, PANEL, corner=True)
    tb(slide, l, t + Inches(0.2), w, Inches(0.7),
       value, size=34, color=color, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(slide, l, t + Inches(0.95), w, Inches(0.4),
       label, size=12, color=MUTED, align=PP_ALIGN.CENTER)

TOTAL = 14
slides = []

# ==================== 1. TITLE ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
bg(s, BG)
# decorative accents
band(s, 0, 0, prs.slide_width, Inches(0.12), CORAL)
band(s, 0, prs.slide_height - Inches(0.12), prs.slide_width, Inches(0.12), TEAL)
# big coral square accent
sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.1), Inches(0.18), Inches(0.6))
sq.fill.solid(); sq.fill.fore_color.rgb = CORAL; sq.line.fill.background()

tb(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.5),
   "PROJECT SHOWCASE", size=14, color=CORAL, bold=True)
tb(s, Inches(0.7), Inches(1.5), Inches(12), Inches(1.4),
   "Smart Attendance &", size=58, color=WHITE, bold=True)
tb(s, Inches(0.7), Inches(2.5), Inches(12), Inches(1.4),
   "Performance Analyzer", size=58, color=CORAL, bold=True)

# tagline
tb(s, Inches(0.7), Inches(3.95), Inches(12), Inches(0.5),
   "A role-based web dashboard system for BTech Computer Science",
   size=20, color=TEXT)
tb(s, Inches(0.7), Inches(4.45), Inches(12), Inches(0.5),
   "Admin  •  5 Teachers  •  10 Students  •  5 CS Courses  •  Zero Backend",
   size=16, color=MUTED)

# stack pills
labels = [("HTML5", CORAL), ("CSS3", TEAL), ("JavaScript", MINT), ("Chart.js", AMBER), ("localStorage", CRIMSON)]
x = Inches(0.7)
for name, c in labels:
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.4), Inches(2.0), Inches(0.55))
    pill.fill.solid(); pill.fill.fore_color.rgb = c
    pill.line.fill.background()
    tf = pill.text_frame
    p = tf.paragraphs[0]; p.text = name
    p.font.size = Pt(16); p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    x += Inches(2.2)

tb(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
   "v1.0  •  Pure client-side  •  No build step  •  Runs in any modern browser",
   size=13, color=DIM)
page_no(s, 1, TOTAL)

# ==================== 2. WHAT IS SAPA? ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "What is SAPA?", "A focused look at the project goals and audience", n=2)

card(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(2.4),
     "Purpose",
     ["Digitize attendance for BTech CS programs",
      "Replace paper registers with real-time tracking",
      "Surface trends through charts and progress bars",
      "Give each role exactly the tools it needs — nothing more"],
     accent=CORAL)

card(s, Inches(6.9), Inches(1.6), Inches(5.8), Inches(2.4),
     "Who uses it",
     ["Admins — manage users, courses, system config",
      "Teachers — mark attendance, view reports",
      "Students — track personal attendance per subject"],
     accent=TEAL)

card(s, Inches(0.7), Inches(4.2), Inches(6.0), Inches(2.7),
     "Why it stands out",
     ["Runs entirely in the browser — no install, no server",
      "Three role-specific dashboards in one SPA shell",
      "Calendar-based date pickers for day / month / year filters",
      "Sample data pre-seeded so the demo works on first load"],
     accent=MINT)

card(s, Inches(6.9), Inches(4.2), Inches(5.8), Inches(2.7),
     "Scope of this build",
     ["5 pre-loaded CS courses (DS, DBMS, OS, CN, SE)",
      "10 students, 5 teacher accounts, 1 admin",
      "750+ attendance records spanning Mar–Jun 2026",
      "Future modules flagged: QR attendance, performance analyzer"],
     accent=AMBER)

# ==================== 3. PROJECT STRUCTURE ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Project Structure", "Files, roles, and what each one owns", n=3)

# Left: file tree
panel(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(5.4), PANEL, corner=True)
tb(s, Inches(0.95), Inches(1.75), Inches(5.6), Inches(0.4),
   "Repository Layout", size=18, color=WHITE, bold=True)

tree = [
    ("Project/", CORAL, True),
    ("├── index.html", WHITE, False),
    ("├── dashboard.html", WHITE, False),
    ("├── css/", AMBER, True),
    ("│   └── style.css", TEXT, False),
    ("├── js/", AMBER, True),
    ("│   ├── app.js (login)", TEXT, False),
    ("│   ├── data.js (storage)", TEXT, False),
    ("│   └── dashboard.js (app)", TEXT, False),
    ("├── objective.md", MUTED, False),
    ("├── suggestion.md", MUTED, False),
    ("├── status.md", MUTED, False),
    ("└── build_ppt.py", MUTED, False),
]
y = Inches(2.2)
for line, color, bold in tree:
    tb(s, Inches(1.1), y, Inches(5.4), Inches(0.32),
       line, size=14, color=color, bold=bold)
    y += Inches(0.34)

# Right: line counts
card(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(2.6),
     "Code footprint",
     ["HTML:    2 pages (login + dashboard shell)",
      "CSS:     728 lines, fully themed",
      "JS:      ~1,400 lines across 3 files",
      "Docs:    objective.md, suggestion.md, status.md",
      "Deps:    1 — Chart.js via CDN"],
     accent=CORAL)

card(s, Inches(7.0), Inches(4.3), Inches(5.7), Inches(2.7),
     "Key facts",
     ["Three role-aware menus, one shared SPA shell",
      "All state lives in localStorage under sarpa_data",
      "Charts re-render on every navigation",
      "Dark / light theme saved to localStorage",
      "Mobile breakpoint at 768px"],
     accent=TEAL)

# ==================== 4. TECH STACK ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Technology Stack", "Every tool used in the build", n=4)

stack = [
    ("HTML5", "Semantic markup for login and SPA shell.\nModal system, nav, <canvas> targets.", CORAL),
    ("CSS3", "Custom-property theming, Grid + Flexbox.\nDark/light mode, 728 lines of styling.", TEAL),
    ("JavaScript", "Vanilla ES6+, IIFE modules.\nEvent-driven SPA, no framework.", MINT),
    ("Chart.js 4.4", "Loaded via CDN.\nBar, line, pie, doughnut charts.", AMBER),
    ("localStorage", "Browser-native persistence.\nJSON serialization, schema v3.", CRIMSON),
    ("Inter Font", "Google Fonts CDN.\nClean legibility across the UI.", TEAL),
]
W, H = Inches(4.0), Inches(2.5)
for i, (name, desc, color) in enumerate(stack):
    col, row = i % 3, i // 3
    l = Inches(0.7) + col * (W + Inches(0.2))
    t = Inches(1.7) + row * (H + Inches(0.2))
    panel(s, l, t, W, H, PANEL, corner=True)
    # header strip
    band(s, l, t, W, Inches(0.55), color)
    tb(s, l, t, W, Inches(0.55), name, size=18, color=WHITE, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, l + Inches(0.2), t + Inches(0.75), W - Inches(0.4), H - Inches(0.85),
       desc, size=13, color=TEXT)

# ==================== 5. ARCHITECTURE ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Architecture", "How a click becomes a chart on screen", n=5)

stages = [
    ("Browser", "User opens index.html.\nLogin form rendered.", CORAL),
    ("Authenticate", "DataStore.authenticate()\nchecks credentials.", TEAL),
    ("Persist", "User saved to\nsarpa_current_user.", MINT),
    ("Redirect", "Loads dashboard.html.\nApp.init() reads user.", AMBER),
    ("Render", "Sidebar + Overview panel\nappear in <1 second.", CRIMSON),
]
W, H = Inches(2.4), Inches(2.0)
gap = Inches(0.15)
total_w = W * 5 + gap * 4
start_x = (prs.slide_width - total_w) / 2
y = Inches(1.9)
for i, (head, body, color) in enumerate(stages):
    l = start_x + i * (W + gap)
    panel(s, l, y, W, H, PANEL, corner=True)
    band(s, l, y, W, Inches(0.5), color)
    tb(s, l, y, W, Inches(0.5), f"Step {i+1}", size=14, color=WHITE, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, l, y + Inches(0.6), W, Inches(0.45),
       head, size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    tb(s, l + Inches(0.15), y + Inches(1.05), W - Inches(0.3), H - Inches(1.1),
       body, size=12, color=TEXT, align=PP_ALIGN.CENTER)

# layer bar at bottom
tb(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.4),
   "Four cooperating layers", size=20, color=CORAL, bold=True, align=PP_ALIGN.CENTER)

layers = [
    ("Presentation", "HTML5 + CSS3", CORAL),
    ("Logic", "JavaScript ES6+", TEAL),
    ("Data", "DataStore + localStorage", MINT),
    ("Visualization", "Chart.js", AMBER),
]
LW, LH = Inches(3.0), Inches(1.4)
ltotal = LW * 4 + Inches(0.3) * 3
lstart = (prs.slide_width - ltotal) / 2
ly = Inches(5.2)
for i, (name, tech, color) in enumerate(layers):
    l = lstart + i * (LW + Inches(0.3))
    panel(s, l, ly, LW, LH, PANEL, corner=True)
    band(s, l, ly, LW, Inches(0.45), color)
    tb(s, l, ly, LW, Inches(0.45), name, size=16, color=WHITE, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, l, ly + Inches(0.55), LW, Inches(0.7),
       tech, size=14, color=TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ==================== 6. DATA MODEL ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Data Model", "Everything stored under localStorage key 'sarpa_data'", n=6)

schema = [
    ("admin",         "{ username, password, name, email }", CORAL),
    ("teachers[]",    "{ id, username, password, name, course, courseCode, email }", TEAL),
    ("students[]",    "{ id, rollno, name, email, batch, courses[], password }", MINT),
    ("attendance[]",  "{ id, studentId, teacherId, courseCode, date, status, remarks }", AMBER),
    ("announcements[]","{ id, text, date, author }", CRIMSON),
    ("auditLog[]",    "{ timestamp, user, action }", CORAL),
    ("config",        "{ minAttendance, lateGraceMinutes, semester dates }", TEAL),
    ("version",       "3  (auto-migrated on schema change)", MUTED),
]
y = Inches(1.7)
for name, shape, color in schema:
    panel(s, Inches(0.7), y, Inches(12.0), Inches(0.55), PANEL, corner=True)
    band(s, Inches(0.7), y, Inches(0.18), Inches(0.55), color)
    tb(s, Inches(1.0), y, Inches(2.5), Inches(0.55),
       name, size=14, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(3.6), y, Inches(8.8), Inches(0.55),
       shape, size=12, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.62)

# ==================== 7. LOGIN & ROLES ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Login & Roles", "Unified login page with role selector", n=7)

card(s, Inches(0.7), Inches(1.6), Inches(5.7), Inches(2.6),
     "Login flow",
     ["User opens index.html",
      "Picks a role from the dropdown (Admin / Teacher / Student)",
      "Enters username and password (default: 123456)",
      "app.js calls DataStore.authenticate()",
      "On success, user is saved to sarpa_current_user",
      "Browser redirects to dashboard.html"],
     accent=CORAL)

card(s, Inches(6.6), Inches(1.6), Inches(6.1), Inches(2.6),
     "Admin credentials",
     ["Username:  admin",
      "Password:  123456",
      "Scope:     full system control"],
     accent=TEAL)

card(s, Inches(0.7), Inches(4.3), Inches(5.7), Inches(2.8),
     "Teacher logins (one per course)",
     ["ds     — Data Structures          (CS201)",
      "dbms   — Database Management      (CS202)",
      "os     — Operating Systems         (CS203)",
      "cn     — Computer Networks         (CS204)",
      "se     — Software Engineering      (CS205)"],
     accent=MINT)

card(s, Inches(6.6), Inches(4.3), Inches(6.1), Inches(2.8),
     "Student logins",
     ["Roll numbers CS2024001 → CS2024010",
      "All students share password 123456",
      "Each student is enrolled in all 5 courses",
      "Personal data is scoped to the logged-in user"],
     accent=AMBER)

# ==================== 8. ADMIN DASHBOARD ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Admin Dashboard", "Full control: users, courses, announcements, audit", n=8)

card(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(2.5),
     "Overview",
     ["Stat cards: teachers, students, today's marks, announcements",
      "Bar chart: teacher distribution across 5 CS courses",
      "Doughnut chart: system-wide statistics at a glance"],
     accent=CORAL)

card(s, Inches(6.9), Inches(1.6), Inches(5.8), Inches(2.5),
     "Manage Teachers",
     ["Full CRUD on teacher records",
      "Assign a teacher to one of 5 CS courses",
      "Cascade delete cleans up attendance",
      "Searchable teacher table"],
     accent=TEAL)

card(s, Inches(0.7), Inches(4.2), Inches(6.0), Inches(2.7),
     "Manage Students",
     ["Add / edit / delete students",
      "Multi-course enrollment via course picker",
      "Fields: name, roll, email, batch, courses[]",
      "Cascade delete on student removal"],
     accent=MINT)

card(s, Inches(6.9), Inches(4.2), Inches(5.8), Inches(2.7),
     "System Tools",
     ["Post announcements with author + timestamp",
      "Audit log of every data change",
      "Configurable min attendance % and late grace",
      "One-click data reset from Settings"],
     accent=AMBER)

# ==================== 9. TEACHER DASHBOARD ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Teacher Dashboard", "Course roster, attendance, and student analytics", n=9)

card(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(2.5),
     "My Course",
     ["View all enrolled students in assigned course",
      "Enroll / remove students with dropdown picker",
      "Student count badge in card header"],
     accent=CORAL)

card(s, Inches(6.9), Inches(1.6), Inches(5.8), Inches(2.5),
     "Mark Attendance",
     ["Calendar date picker for the session",
      "Present / Absent / Late toggle cards",
      "Mark All Present / Mark All Absent bulk actions",
      "Auto-save to localStorage on every click"],
     accent=TEAL)

card(s, Inches(0.7), Inches(4.2), Inches(6.0), Inches(2.7),
     "View Attendance",
     ["Three filters: By Day, By Month, By Year",
      "Calendar picker for precise date / month range",
      "Color-coded badges: green / red / amber",
      "Tabular output with roll, name, date, status"],
     accent=MINT)

card(s, Inches(6.9), Inches(4.2), Inches(5.8), Inches(2.7),
     "Student Analytics",
     ["Per-student summary table",
      "Columns: Total, Present, Absent, Late, %",
      "Color-coded progress bars by threshold",
      "Spot at-risk students in a glance"],
     accent=AMBER)

# ==================== 10. STUDENT DASHBOARD ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Student Dashboard", "Personal attendance, course breakdown, charts", n=10)

card(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(2.5),
     "My Attendance",
     ["Personal log with day / month / year filters",
      "Calendar picker for precise dates",
      "Color-coded status badges",
      "Course name shown alongside each entry"],
     accent=CORAL)

card(s, Inches(6.9), Inches(1.6), Inches(5.8), Inches(2.5),
     "Course Breakdown",
     ["Overall attendance % with progress bar",
      "Stats per CS course: total, present, absent, late, %",
      "Color-coded bars for instant assessment",
      "Covers DS, DBMS, OS, CN, SE"],
     accent=TEAL)

card(s, Inches(0.7), Inches(4.2), Inches(6.0), Inches(2.7),
     "Charts (4 types)",
     ["Pie: overall present / absent / late split",
      "Line: monthly attendance trend over semester",
      "Bar: course-wise attendance % comparison",
      "Doughnut: per-course proportional breakdown"],
     accent=MINT)

card(s, Inches(6.9), Inches(4.2), Inches(5.8), Inches(2.7),
     "Overview Stats",
     ["Stat cards: overall %, present, absent, courses",
      "Personal pie chart of attendance distribution",
      "Monthly line chart of attendance trend",
      "Auto-calculated across all enrolled courses"],
     accent=AMBER)

# ==================== 11. CHARTS & VISUAL DESIGN ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Charts & Visual Design", "How Chart.js brings the data to life", n=11)

# 4 chart cards with descriptions
charts = [
    ("Bar",  "Monthly + course-wise attendance %", CORAL, "▌▌▌▌"),
    ("Line", "Attendance progression over time",   TEAL, "∿∿∿"),
    ("Pie",  "Present / Absent / Late split",      MINT, "◐◑◒"),
    ("Doughnut", "System stats + per-course split", AMBER, "◎"),
]
CW, CH = Inches(3.0), Inches(2.5)
for i, (name, desc, color, glyph) in enumerate(charts):
    col, row = i % 4, i // 4
    l = Inches(0.7) + col * (CW + Inches(0.15))
    t = Inches(1.7) + row * (CH + Inches(0.2))
    panel(s, l, t, CW, CH, PANEL, corner=True)
    # visual mock
    mock = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.4), t + Inches(0.3), Inches(2.2), Inches(1.1))
    mock.fill.solid(); mock.fill.fore_color.rgb = color
    mock.line.color.rgb = WHITE; mock.line.width = Pt(0.5)
    tf = mock.text_frame
    p = tf.paragraphs[0]; p.text = glyph
    p.font.size = Pt(28); p.font.color.rgb = WHITE; p.font.bold = True
    p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb(s, l, t + Inches(1.55), CW, Inches(0.4),
       name, size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    tb(s, l + Inches(0.15), t + Inches(1.95), CW - Inches(0.3), Inches(0.5),
       desc, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Design system card
card(s, Inches(0.7), Inches(4.6), Inches(12.0), Inches(2.3),
     "Design system",
     ["Palette: deep navy + coral accents + mint / amber / crimson status colors",
      "Typography: Inter via Google Fonts",
      "Theme: dark / light toggle saved to localStorage",
      "Charts dynamically created and destroyed on every page switch",
      "Responsive breakpoint at 768px stacks the sidebar and panels"],
     accent=CORAL, head=20, body=14)

# ==================== 12. KEY FEATURES SUMMARY ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Feature Highlights", "What makes the build complete", n=12)

features = [
    ("Auth", "Role-based login\nfor 3 user types", CORAL),
    ("CRUD", "Admin can add,\nedit, delete users", TEAL),
    ("Calendar", "Native date pickers\nfor all filters", MINT),
    ("Progress", "Color-coded bars\nfor attendance %", AMBER),
    ("Audit", "Every change\nis logged", CRIMSON),
    ("Responsive", "Mobile-friendly\nbelow 768px", CORAL),
    ("Theme", "Dark / light\nmode toggle", TEAL),
    ("Seeded Data", "750+ records\npre-loaded", MINT),
]
W, H = Inches(3.0), Inches(1.4)
for i, (name, desc, color) in enumerate(features):
    col, row = i % 4, i // 4
    l = Inches(0.7) + col * (W + Inches(0.15))
    t = Inches(1.7) + row * (H + Inches(0.2))
    panel(s, l, t, W, H, PANEL, corner=True)
    band(s, l, t, Inches(0.18), H, color)
    tb(s, l + Inches(0.3), t + Inches(0.15), W - Inches(0.4), Inches(0.4),
       name, size=16, color=color, bold=True)
    tb(s, l + Inches(0.3), t + Inches(0.55), W - Inches(0.4), H - Inches(0.6),
       desc, size=12, color=TEXT)

# planned features banner
band(s, Inches(0.7), Inches(5.0), Inches(12.0), Inches(0.6), TEAL)
tb(s, Inches(0.9), Inches(5.0), Inches(11.6), Inches(0.6),
   "PLANNED  •  QR attendance  •  Performance analyzer  •  Leave requests  •  PDF / Excel exports  •  Email alerts",
   size=14, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

card(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(1.2),
     "Why this matters",
     ["Self-contained demo — no server, no setup, no signup. Open the HTML and it works."],
     accent=CORAL, head=16, body=13)

# ==================== 13. KEY METRICS ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
chrome(s, "Project Metrics", "The build at a glance", n=13)

metrics = [
    ("3", "Role-based\ndashboards", CORAL),
    ("5", "BTech CS\ncourses", TEAL),
    ("10", "Pre-loaded\nstudents", MINT),
    ("750+", "Sample\nattendance rows", AMBER),
    ("4", "Chart.js\nchart types", CRIMSON),
    ("~2,200", "Total lines\nof code", CORAL),
    ("1", "External\ndependency", TEAL),
    ("0", "Backend\nservices", MINT),
]
W, H = Inches(3.0), Inches(2.0)
for i, (val, label, color) in enumerate(metrics):
    col, row = i % 4, i // 4
    l = Inches(0.7) + col * (W + Inches(0.15))
    t = Inches(1.7) + row * (H + Inches(0.2))
    panel(s, l, t, W, H, PANEL, corner=True)
    tb(s, l, t + Inches(0.2), W, Inches(0.9),
       val, size=46, color=color, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, l, t + Inches(1.15), W, Inches(0.7),
       label.replace("\n", "  •  "), size=13, color=MUTED, align=PP_ALIGN.CENTER)

# ==================== 14. THANK YOU ====================
s = prs.slides.add_slide(BLANK); slides.append(s)
bg(s, BG)
band(s, 0, 0, prs.slide_width, Inches(0.12), CORAL)
band(s, 0, prs.slide_height - Inches(0.12), prs.slide_width, Inches(0.12), TEAL)

sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(2.2), Inches(0.18), Inches(0.6))
sq.fill.solid(); sq.fill.fore_color.rgb = CORAL; sq.line.fill.background()

tb(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
   "END OF SHOWCASE", size=14, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
tb(s, Inches(0.7), Inches(2.6), Inches(12), Inches(1.4),
   "Thank You", size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb(s, Inches(0.7), Inches(4.0), Inches(12), Inches(0.5),
   "Questions, feedback, and ideas are welcome.",
   size=20, color=MUTED, align=PP_ALIGN.CENTER)

# quick recap pills
labels = [("Admin: admin / 123456", CORAL),
          ("Teachers: ds · dbms · os · cn · se", TEAL),
          ("Students: CS2024001–CS2024010", MINT),
          ("Default password: 123456", AMBER)]
x = Inches(0.7)
for label, color in labels:
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.1), Inches(3.0), Inches(0.6))
    pill.fill.solid(); pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    tf = pill.text_frame
    p = tf.paragraphs[0]; p.text = label
    p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    x += Inches(3.1)

tb(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
   "SAPA  •  Smart Attendance & Performance Analyzer  •  v1.0",
   size=14, color=DIM, align=PP_ALIGN.CENTER)
page_no(s, 14, TOTAL)

# ---------- Save ----------
out = r"C:\Users\ulhas\OneDrive\Documents\Project\SAPA_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
